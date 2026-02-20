"""
Generate final report PPT (Samsung Electronics/Display standard format).

Derived from final paper (FPCB_Crack_Detection_Final_Report.docx).
- Samsung theme: #1428A0 blue, Arial, min 16pt body
- VS comparison slides (side-by-side)
- Section divider slides
- All text in English

Output: reports/deliverables/FPCB_Crack_Detection_Final_Report.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

repo_root = Path(__file__).resolve().parent.parent
reports_dir = repo_root / "reports"
analysis_dir = reports_dir / "crack_detection_analysis"
videos_dir = reports_dir / "deliverables" / "videos"

# Samsung brand colors (hex #1428A0 = RGB 20,40,160)
SAMSUNG_BLUE = (20, 40, 160)
SAMSUNG_BLUE_LIGHT = (80, 110, 200)
WHITE = (255, 255, 255)
DARK_GRAY = (50, 50, 50)
HIGHLIGHT_GREEN = (0, 150, 80)

# Slide dimensions and margins (10 x 7.5 inches)
SLIDE_W = 10.0
SLIDE_H = 7.5
MARGIN = 0.4
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_H = SLIDE_H - 2 * MARGIN - 0.6  # Leave room for title


def _inches(v: float) -> float:
    return v


def _apply_samsung_theme(slide, use_blue_bg: bool = False) -> None:
    """Apply Samsung color theme to slide."""
    from pptx.dml.color import RGBColor

    if use_blue_bg:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*SAMSUNG_BLUE)
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*WHITE)


def _set_font(shape, font_name: str = "Arial", font_size: int = 12, bold: bool = False, color_rgb: tuple = DARK_GRAY) -> None:
    """Set font for text in shape."""
    from pptx.dml.color import RGBColor

    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = _pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(*color_rgb)


def _pt(size: int):
    from pptx.util import Pt
    return Pt(size)


def _add_title_slide(prs, title: str, subtitle: str) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide, use_blue_bg=True)

    left, top, width, height = Inches(MARGIN), Inches(2.5), Inches(CONTENT_W), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = 1  # Center
    _set_font(tx, font_size=36, bold=True, color_rgb=WHITE)

    left, top, width, height = Inches(MARGIN), Inches(4), Inches(CONTENT_W), Inches(0.8)
    tx2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.alignment = 1
    _set_font(tx2, font_size=18, color_rgb=WHITE)


def _add_section_divider_slide(prs, section_title: str) -> None:
    """Full-slide section divider (blue bg)."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide, use_blue_bg=True)
    left, top, width, height = Inches(MARGIN), Inches(3.0), Inches(CONTENT_W), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.alignment = 1
    _set_font(tx, font_size=36, bold=True, color_rgb=WHITE)


def _add_section_slide(prs, title: str, bullets: list[str], blue_header: bool = True) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.9)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _set_font(tx, font_size=26, bold=True, color_rgb=SAMSUNG_BLUE)

    top = Inches(1.3)
    for bullet in bullets:
        left, width, height = Inches(MARGIN), Inches(CONTENT_W), Inches(0.55)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"•  {bullet}"
        _set_font(tx, font_size=16, color_rgb=DARK_GRAY)
        top += Inches(0.48)


def _add_table_slide(prs, title: str, headers: list[str], rows: list[list], highlight_cols: list[int] = None) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.6)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _set_font(tx, font_size=24, bold=True, color_rgb=SAMSUNG_BLUE)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_left, table_top = Inches(MARGIN), Inches(1.0)
    table_width = Inches(CONTENT_W)
    row_height = Inches(0.45)
    table_height = row_height * n_rows
    table = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height).table

    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = str(h)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(*SAMSUNG_BLUE)
        p = c.text_frame.paragraphs[0]
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)
        p.alignment = 1

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j < n_cols:
                c = table.cell(i + 1, j)
                c.text = str(val)
                is_highlight = highlight_cols and j in highlight_cols
                if is_highlight:
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(*HIGHLIGHT_GREEN)
                    c.text_frame.paragraphs[0].font.color.rgb = RGBColor(*WHITE)
                p = c.text_frame.paragraphs[0]
                p.font.name = "Arial"
                p.font.size = Pt(11)
                p.alignment = 1


def _add_image_slide(prs, title: str, img_path: Path, full_size: bool = True) -> None:
    from pptx.util import Inches

    if not img_path.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _set_font(tx, font_size=22, bold=True, color_rgb=SAMSUNG_BLUE)

    img_left, img_top = Inches(MARGIN), Inches(1.0)
    img_w = Inches(CONTENT_W)
    img_h = Inches(5.5) if full_size else Inches(4.5)
    slide.shapes.add_picture(str(img_path), img_left, img_top, width=img_w, height=img_h)


def _add_vs_comparison_slide(prs, title: str, left_img: Path, right_img: Path, left_label: str, right_label: str, highlight_text: str = "") -> None:
    """Side-by-side comparison slide with VS layout."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if not left_img.exists() or not right_img.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = 1
    _set_font(tx, font_size=24, bold=True, color_rgb=SAMSUNG_BLUE)

    # Left image
    img_w = Inches(4.0)
    img_h = Inches(4.5)
    slide.shapes.add_picture(str(left_img), Inches(MARGIN), Inches(1.0), width=img_w, height=img_h)
    lb = slide.shapes.add_textbox(Inches(MARGIN), Inches(5.6), img_w, Inches(0.4))
    lb.text_frame.paragraphs[0].text = left_label
    lb.text_frame.paragraphs[0].alignment = 1
    _set_font(lb, font_size=14, bold=True, color_rgb=SAMSUNG_BLUE)

    # VS label
    vs = slide.shapes.add_textbox(Inches(4.6), Inches(3.0), Inches(0.8), Inches(0.6))
    vs.text_frame.paragraphs[0].text = "VS"
    vs.text_frame.paragraphs[0].alignment = 1
    _set_font(vs, font_size=28, bold=True, color_rgb=SAMSUNG_BLUE)

    # Right image
    right_left = Inches(SLIDE_W - MARGIN - 4.0)
    slide.shapes.add_picture(str(right_img), right_left, Inches(1.0), width=img_w, height=img_h)
    rb = slide.shapes.add_textbox(right_left, Inches(5.6), img_w, Inches(0.4))
    rb.text_frame.paragraphs[0].text = right_label
    rb.text_frame.paragraphs[0].alignment = 1
    _set_font(rb, font_size=14, bold=True, color_rgb=SAMSUNG_BLUE)

    if highlight_text:
        hl = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.2), Inches(CONTENT_W), Inches(0.5))
        hl.text_frame.paragraphs[0].text = f"Key difference: {highlight_text}"
        hl.text_frame.paragraphs[0].alignment = 1
        _set_font(hl, font_size=12, color_rgb=HIGHLIGHT_GREEN)


def _add_three_way_comparison(prs, title: str, img_paths: list[Path], labels: list[str]) -> None:
    """Three images in a row for DREAM vs PatchCore vs Ensemble."""
    from pptx.util import Inches

    if len(img_paths) < 3 or len(labels) < 3:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = 1
    _set_font(tx, font_size=22, bold=True, color_rgb=SAMSUNG_BLUE)

    img_w = Inches(2.9)
    img_h = Inches(4.2)
    gap = Inches(0.15)
    total_w = 3 * 2.9 + 2 * 0.15
    start_x = (SLIDE_W - total_w) / 2

    for i, (path, label) in enumerate(zip(img_paths[:3], labels[:3])):
        if not path.exists():
            continue
        x = Inches(start_x + i * (2.9 + 0.15))
        slide.shapes.add_picture(str(path), x, Inches(1.0), width=img_w, height=img_h)
        lb = slide.shapes.add_textbox(x, Inches(5.3), img_w, Inches(0.4))
        lb.text_frame.paragraphs[0].text = label
        lb.text_frame.paragraphs[0].alignment = 1
        _set_font(lb, font_size=11, bold=True, color_rgb=SAMSUNG_BLUE)


def _add_video_slide(prs, title: str, video_path: Path, poster_path: Path | None = None) -> bool:
    from pptx.util import Inches

    if not video_path.exists():
        return False
    poster = poster_path if poster_path and poster_path.exists() else None
    if not poster and (analysis_dir / "confusion_matrix_ensemble.png").exists():
        poster = analysis_dir / "confusion_matrix_ensemble.png"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_samsung_theme(slide)

    left, top, width, height = Inches(MARGIN), Inches(MARGIN), Inches(CONTENT_W), Inches(0.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _set_font(tx, font_size=22, bold=True, color_rgb=SAMSUNG_BLUE)

    vid_w, vid_h = Inches(CONTENT_W), Inches(5.5)
    try:
        slide.shapes.add_movie(
            str(video_path),
            Inches(MARGIN), Inches(1.0),
            vid_w, vid_h,
            poster_frame_image=str(poster) if poster else None,
            mime_type="video/mp4",
        )
        return True
    except Exception:
        if poster:
            slide.shapes.add_picture(str(poster), Inches(MARGIN), Inches(1.0), width=vid_w, height=vid_h)
        tx2 = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.6), Inches(CONTENT_W), Inches(0.4))
        tx2.text_frame.paragraphs[0].text = f"[Video: {video_path.name}] Double-click to play"
        _set_font(tx2, font_size=11, color_rgb=DARK_GRAY)
        return True


def main() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    normal_map = analysis_dir / "vector_map_normal.png"
    crack_map = analysis_dir / "vector_map_crack.png"
    dream_img = analysis_dir / "confusion_matrix_dream.png"
    patch_img = analysis_dir / "confusion_matrix_patchcore.png"
    ensemble_img = analysis_dir / "confusion_matrix_ensemble.png"
    insights_img = analysis_dir / "insights_summary.png"
    v1 = videos_dir / "02_analysis_process_log.mp4"
    v2 = videos_dir / "01_vector_map_visualization.mp4"
    v3 = videos_dir / "03_confusion_matrix_results.mp4"

    # ─── IMRaD Structure: Introduction, Methods, Results, Discussion, Appendix ───

    # 1. Title
    _add_title_slide(prs, "FPCB Bending Process Crack Detection AI", "Final Report | IMRaD Format | 2026-02-20")

    # 2. Table of Contents (IMRaD)
    _add_section_slide(prs, "Table of Contents", [
        "1. Introduction — Background, Problem, Objectives, Related Work",
        "2. Methods — Dataset, Synthetic Construction, Models, Ensemble",
        "3. Results — Performance, Comparisons, Figures",
        "4. Discussion — Conclusions, Recommendations, Future Work",
        "Appendix — Definitions, Dataset",
    ])

    # ═══════════════════════════════════════════════════════════════════════════════
    # INTRODUCTION
    # ═══════════════════════════════════════════════════════════════════════════════

    _add_section_slide(prs, "1. Introduction — Background", [
        "FPCB: layered organic film + copper traces; used in foldable displays.",
        "Copper cracks occur during bending; early detection critical for quality control.",
        "False alarms (FP) cause unnecessary line stops; reduce throughput.",
        "Traditional inspection struggles with subtle anomalies and illumination variation.",
    ])

    _add_section_slide(prs, "1. Introduction — Problem & Objectives", [
        "Problem: Baseline 86–88% precision, 93–130 FPs. Primary FP source: light_distortion.",
        "Objectives: (1) Precision ≥99%, (2) FP near-zero, (3) Robustness to light_distortion.",
        "Scope: Physics-informed synthetic motion; DREAM + PatchCore ensemble.",
    ])

    _add_section_slide(prs, "1. Introduction — Related Work & Our Differentiation", [
        "DRAEM [1]: Image-level; 98.1% ROC AUC on MVTec. PatchCore [2]: 99.6% AUROC.",
        "Our work: Motion features (velocity, curvature); temporal focus; DREAM ∧ PatchCore ensemble.",
        "Logical-AND ensemble achieves 100% precision, FP=0.",
    ])

    # ═══════════════════════════════════════════════════════════════════════════════
    # METHODS
    # ═══════════════════════════════════════════════════════════════════════════════

    _add_section_divider_slide(prs, "2. Methods")

    _add_section_slide(prs, "2. Methods — Overview", [
        "Dataset: 7 scenarios (normal, crack, light_distortion, micro_crack, etc.).",
        "Models: DREAM (DRAEM) + PatchCore; Ensemble = DREAM ∧ PatchCore (both agree).",
        "4 validation loops: augmentation, threshold, ensemble, MIN_PRECISION 0.997.",
    ])

    _add_table_slide(prs, "2. Methods — Synthetic Dataset Composition", ["Scenario", "Count", "Label"], [
        ["normal", "1,000", "0 (normal)"],
        ["light_distortion", "50", "0 (normal)"],
        ["crack", "50", "1 (crack)"],
        ["uv_overcured", "30", "1"],
        ["micro_crack", "10", "1"],
        ["pre_damaged", "20", "1"],
        ["thick_panel", "20", "0"],
    ])

    _add_section_slide(prs, "2. Methods — Synthetic Data Trustworthiness", [
        "Physics alignment: Bend progression (straight→arc→U); shockwave at crack frame.",
        "Scenario semantics: light_distortion, micro_crack map to real conditions.",
        "Reproducibility: Explicit seed, ScenarioParams; deterministic regeneration.",
    ])

    _add_section_slide(prs, "2. Methods — Model Architecture", [
        "DREAM (DRAEM [1]): Reconstruction + discriminative head; temporal patterns.",
        "PatchCore [2]: Memory bank + k-NN; feature-space distance.",
        "Ensemble: Logical AND — both must predict Crack; mutual FP filtering.",
    ])

    if v1.exists():
        _add_video_slide(prs, "2. Methods — [Video] Analysis Process Log", v1, insights_img)

    # ═══════════════════════════════════════════════════════════════════════════════
    # RESULTS
    # ═══════════════════════════════════════════════════════════════════════════════

    _add_section_divider_slide(prs, "3. Results")

    _add_section_slide(prs, "3. Results — Overview", [
        "Ensemble achieved Precision 100% and FP = 0.",
        "light_distortion (8/8) and micro_crack (2/2) correctly classified.",
        "4 development-validation loops led to incremental improvement.",
    ])

    _add_table_slide(prs, "3. Results — Performance Overview", ["Metric", "Baseline", "Final (Ensemble)", "Improvement"], [
        ["Precision", "86–88%", "100%", "+12–14%p"],
        ["False Positive", "93–130", "0", "100% reduction"],
        ["light_distortion (normal)", "0%", "100%", "+100%p"],
    ], highlight_cols=[2, 3])

    _add_table_slide(prs, "3. Results — Model Comparison", ["Model", "Precision", "FP", "Recall"], [
        ["DREAM", "99.83%", "1", "67.8%"],
        ["PatchCore", "99.82%", "1", "65.5%"],
        ["Ensemble (DREAM ∧ PatchCore)", "100%", "0", "65.2%"],
    ], highlight_cols=[1, 2])

    _add_table_slide(prs, "3. Results — Hard Subset", ["Scenario", "Result", "Note"], [
        ["light_distortion", "8/8 (100%)", "Correctly normal"],
        ["micro_crack", "2/2 (100%)", "Correctly crack"],
    ])

    if normal_map.exists() and crack_map.exists():
        _add_vs_comparison_slide(
            prs,
            "3. Results — Vector Map: Normal vs Crack",
            normal_map, crack_map,
            "Normal: Smooth velocity/acceleration",
            "Crack: Shockwave & vibration at crack frame",
            "Crack shows acceleration spike; Normal shows smooth progression",
        )

    if v2.exists():
        _add_video_slide(prs, "3. Results — [Video] Vector Map Visualization", v2, normal_map)

    if dream_img.exists() and patch_img.exists() and ensemble_img.exists():
        _add_three_way_comparison(
            prs,
            "3. Results — Confusion Matrix: DREAM vs PatchCore vs Ensemble",
            [dream_img, patch_img, ensemble_img],
            ["DREAM (FP=1)", "PatchCore (FP=1)", "Ensemble (FP=0)"],
        )

    if ensemble_img.exists():
        _add_image_slide(prs, "3. Results — Ensemble Confusion Matrix (Final, Precision 100%)", ensemble_img)

    if v3.exists():
        _add_video_slide(prs, "3. Results — [Video] Confusion Matrix Results", v3, ensemble_img)

    _add_table_slide(prs, "3. Results — Development-Validation Loop", ["Loop", "Changes", "Precision", "FP"], [
        ["Baseline", "-", "86–88%", "93–130"],
        ["1", "light_distortion 50, thick_panel", "99.15%", "5"],
        ["2", "light_distortion augmentation", "99.13%", "5"],
        ["3", "Ensemble (both agree)", "99.65%", "2"],
        ["4", "MIN_PRECISION 0.997", "100%", "0"],
    ], highlight_cols=[2, 3])

    # ═══════════════════════════════════════════════════════════════════════════════
    # DISCUSSION
    # ═══════════════════════════════════════════════════════════════════════════════

    _add_section_divider_slide(prs, "4. Discussion")

    _add_section_slide(prs, "4. Discussion — Conclusions", [
        "Precision 100% and FP=0 achieved; false alarm goal met.",
        "Ensemble provides mutual FP filtering; each base model had FP=1.",
        "light_distortion 50 samples → 100% correct; Recall 65% acceptable.",
    ])

    _add_section_slide(prs, "4. Discussion — Recommendations & Future Work", [
        "Re-validate with real FPCB footage; thresholds provisional until calibration.",
        "Future: Mixed training (synthetic + few real defects); incremental integration.",
        "Threshold recalibration on real val set; A/B testing and drift monitoring.",
    ])

    # ═══════════════════════════════════════════════════════════════════════════════
    # APPENDIX
    # ═══════════════════════════════════════════════════════════════════════════════

    _add_section_slide(prs, "Appendix — Confusion Matrix Definitions", [
        "TN: Correctly normal. FP: Normal→crack (target). FN: Crack→normal. TP: Correctly crack.",
        "Precision = TP / (TP + FP) → 100% when FP=0.",
    ])

    if insights_img.exists():
        _add_image_slide(prs, "Appendix — Performance Summary (Insights)", insights_img)

    out_path = reports_dir / "deliverables" / "FPCB_Crack_Detection_Final_Report.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
